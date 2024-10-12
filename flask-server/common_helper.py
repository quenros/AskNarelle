import fitz  # PyMuPDF
import docx
import io
from pptx import Presentation

def read_pdf(blob_data):
    with io.BytesIO(blob_data) as temp_file:
        document = fitz.open("pdf", temp_file)
        text = ""
        for page_num in range(len(document)):
            page = document.load_page(page_num)
            text += page.get_text()
    return text

def read_docx(blob_data):
    with io.BytesIO(blob_data) as temp_file:
        doc = docx.Document(temp_file)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

def read_txt(blob_data):
    with io.BytesIO(blob_data) as temp_file:
        text = temp_file.read().decode('utf-8') 
    return text

def read_pptx(blob_data):
    with io.BytesIO(blob_data) as temp_file:
        presentation = Presentation(temp_file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text